"""
Fox command pattern — one ToolCommand subclass per tool, CommandRegistry with
cache-before-execute via Storage.
"""

import hashlib
import json
import os
import re
import subprocess
import sys
import time
from abc import ABC, abstractmethod
from dataclasses import dataclass
from typing import Optional, TYPE_CHECKING

if TYPE_CHECKING:
    from src.storage import Storage


# ── Result ────────────────────────────────────────────────────────────────────

@dataclass
class CommandResult:
    output: str
    success: bool
    elapsed: float
    exit_code: Optional[int] = None


# ── Base ──────────────────────────────────────────────────────────────────────

class ToolCommand(ABC):
    name: str = ""

    def __init__(self, args: dict):
        self.args: dict = args
        self.timestamp: float = 0.0
        self.result: Optional[CommandResult] = None

    @abstractmethod
    def execute(self) -> CommandResult: ...

    def undo(self) -> Optional[str]:
        return None

    def args_hash(self) -> str:
        return hashlib.sha256(
            json.dumps(self.args, sort_keys=True).encode()
        ).hexdigest()[:16]

    @property
    def metadata(self) -> dict:
        return {
            "name": self.name,
            "args_summary": self._summarize_args(),
            "timestamp": self.timestamp,
            "success": self.result.success if self.result else None,
        }

    def _summarize_args(self) -> str:
        parts = []
        for k, v in self.args.items():
            s = str(v)
            parts.append(f"{k}={s[:80]}..." if len(s) > 80 else f"{k}={s}")
        return ", ".join(parts)

    @staticmethod
    def _truncate(text: str, limit: int = 20_000) -> str:
        return text[:limit] if len(text) > limit else text


# ── Subclasses ────────────────────────────────────────────────────────────────

class RunBashCommand(ToolCommand):
    name = "run_bash"

    # Pre-execution guards — catch known bad patterns before wasting a shell call
    _BAD_PATTERNS = [
        (
            re.compile(r'\bpython3?\s+-m\s+pptx\b'),
            "Error: `python -m pptx` is not valid — pptx is a library, not a CLI tool. "
            "Write a Python script using `from pptx import Presentation` and run it with "
            "`python3 script.py` or a heredoc."
        ),
        (
            re.compile(r'(?<![3])python\s'),  # bare `python ` without `python3`
            None,  # None = let it run but post-process the "not found" error
        ),
    ]

    def execute(self) -> CommandResult:
        t0 = time.time()
        cmd = self.args["command"]

        # Pre-execution: block known-wrong patterns immediately
        if re.search(r'\bpython3?\s+-m\s+pptx\b', cmd):
            msg = (
                "Error: `python -m pptx` is invalid — pptx is a library, not a CLI tool.\n"
                "Write a Python script and run it:\n"
                "  python3 << 'PYEOF'\n"
                "  from pptx import Presentation\n"
                "  prs = Presentation()\n"
                "  slide = prs.slides.add_slide(prs.slide_layouts[1])\n"
                "  slide.shapes.title.text = 'Title'\n"
                "  slide.placeholders[1].text = 'Body text'\n"
                "  prs.save('/path/to/output.pptx')\n"
                "  print('Saved')\n"
                "  PYEOF"
            )
            self.result = CommandResult(msg, False, time.time() - t0)
            return self.result

        try:
            r = subprocess.run(
                cmd, shell=True,
                capture_output=True, text=True, timeout=120, cwd=os.getcwd(),
            )
            out = r.stdout or ""
            if r.stderr:
                out += ("\n--- stderr ---\n" + r.stderr) if out else r.stderr
            if r.returncode != 0:
                out += f"\n[exit code: {r.returncode}]"
            # Post-process: enrich "python: not found" errors
            if "python: not found" in out or "python: command not found" in out:
                out = "Error: `python` not found — use `python3` instead.\n" + out
            self.result = CommandResult(
                self._truncate(out) or "(no output)",
                r.returncode == 0, time.time() - t0, r.returncode,
            )
        except subprocess.TimeoutExpired:
            self.result = CommandResult("(command timed out)", False, time.time() - t0)
        except Exception as e:
            self.result = CommandResult(f"Error: {e}", False, time.time() - t0)
        return self.result


class RunPythonCommand(ToolCommand):
    name = "run_python"

    def __init__(self, args: dict, work_dir: str):
        super().__init__(args)
        self.work_dir = work_dir

    # Detects large inline string literals that are almost certainly pasted file content
    _INLINE_DATA_RE = re.compile(r"('''|\"\"\")(.{200,}?)\1", re.DOTALL)

    def execute(self) -> CommandResult:
        t0 = time.time()
        script = self.args["script"]
        # Strip markdown code fences that LLMs often include
        script = re.sub(r'^```(?:python)?\s*\n?', '', script)
        script = re.sub(r'\n?```\s*$', '', script)

        # Guard: reject scripts that hardcode large data blobs inline.
        # The model should always read from the file instead.
        if self._INLINE_DATA_RE.search(script):
            msg = (
                "Error: script contains a large inline string literal — "
                "do not hardcode data in scripts. "
                f"Read from the file instead:\n"
                f"  with open('{self.work_dir}/user_input.txt') as f:\n"
                f"      text = f.read()"
            )
            self.result = CommandResult(msg, False, time.time() - t0)
            return self.result
        script_path = os.path.join(self.work_dir, "_script.py")
        try:
            with open(script_path, "w") as f:
                f.write(script)
            r = subprocess.run(
                [sys.executable, script_path],
                capture_output=True, text=True, timeout=120, cwd=os.getcwd(),
                env={**os.environ, "WORK_DIR": self.work_dir},
            )
            out = r.stdout or ""
            if r.stderr:
                out += ("\n--- stderr ---\n" + r.stderr) if out else r.stderr
            if r.returncode != 0:
                out += f"\n[exit code: {r.returncode}]"
            self.result = CommandResult(
                self._truncate(out) or "(no output)",
                r.returncode == 0, time.time() - t0, r.returncode,
            )
        except subprocess.TimeoutExpired:
            self.result = CommandResult("(command timed out)", False, time.time() - t0)
        except Exception as e:
            self.result = CommandResult(f"Error: {e}", False, time.time() - t0)
        return self.result


class ReadFileCommand(ToolCommand):
    name = "read_file"

    def execute(self) -> CommandResult:
        t0 = time.time()
        try:
            path = os.path.expanduser(self.args["path"])
            if not os.path.isabs(path):
                path = os.path.join(os.getcwd(), path)
            with open(path) as f:
                lines = f.readlines()
            start = self.args.get("start_line", 1) - 1
            end = self.args.get("end_line", len(lines))
            selected = lines[max(0, start):end]
            numbered = [
                f"{i + max(0, start) + 1:4d} | {line}"
                for i, line in enumerate(selected)
            ]
            output = self._truncate("".join(numbered)) or "(empty file)"
            self.result = CommandResult(output, True, time.time() - t0)
        except Exception as e:
            self.result = CommandResult(f"Error: {e}", False, time.time() - t0)
        return self.result


_BINARY_EXTENSIONS = {
    ".pptx", ".ppt", ".xlsx", ".xls", ".docx", ".doc",
    ".pdf", ".png", ".jpg", ".jpeg", ".gif", ".zip",
    ".tar", ".gz", ".bz2", ".whl", ".egg",
}


class WriteFileCommand(ToolCommand):
    name = "write_file"

    def __init__(self, args: dict):
        super().__init__(args)
        self._previous_content: Optional[str] = None

    def execute(self) -> CommandResult:
        t0 = time.time()
        try:
            path = os.path.expanduser(self.args["path"])
            if not os.path.isabs(path):
                path = os.path.join(os.getcwd(), path)

            # Block text writes to binary format extensions — these files require
            # library-generated binary content and cannot be created by writing text.
            ext = os.path.splitext(path)[1].lower()
            if ext in _BINARY_EXTENSIONS:
                msg = (
                    f"Error: write_file cannot create {ext} files — they require binary "
                    f"library output, not text. Use run_bash with python-pptx/openpyxl/etc:\n"
                    f"  1. run_bash: pip install python-pptx -q && echo OK\n"
                    f"  2. run_bash: python3 << 'EOF'\n"
                    f"from pptx import Presentation\n"
                    f"prs = Presentation()\n"
                    f"slide = prs.slides.add_slide(prs.slide_layouts[1])\n"
                    f"slide.shapes.title.text = 'Title'\n"
                    f"slide.placeholders[1].text = 'Content'\n"
                    f"prs.save('{path}')\n"
                    f"print('Saved {path}')\n"
                    f"EOF"
                )
                self.result = CommandResult(msg, False, time.time() - t0)
                return self.result

            if os.path.exists(path):
                with open(path) as f:
                    self._previous_content = f.read()
            os.makedirs(os.path.dirname(path) or ".", exist_ok=True)
            with open(path, "w") as f:
                f.write(self.args["content"])
            output = f"Wrote {len(self.args['content'])} bytes to {path}"
            self.result = CommandResult(output, True, time.time() - t0)
        except Exception as e:
            self.result = CommandResult(f"Error: {e}", False, time.time() - t0)
        return self.result

    def undo(self) -> Optional[str]:
        path = os.path.expanduser(self.args["path"])
        if not os.path.isabs(path):
            path = os.path.join(os.getcwd(), path)
        try:
            if self._previous_content is not None:
                with open(path, "w") as f:
                    f.write(self._previous_content)
                return f"Restored {path}"
            elif os.path.exists(path):
                os.remove(path)
                return f"Removed {path}"
        except Exception as e:
            return f"Undo failed: {e}"
        return None


class GrepSearchCommand(ToolCommand):
    name = "grep_search"

    def execute(self) -> CommandResult:
        t0 = time.time()
        try:
            cmd = ["grep", "-rn", "--color=never"]
            if self.args.get("include"):
                cmd += [f"--include={self.args['include']}"]
            cmd.append(self.args["pattern"])
            cmd.append(self.args.get("path", "."))
            r = subprocess.run(
                cmd, capture_output=True, text=True, timeout=30, cwd=os.getcwd()
            )
            self.result = CommandResult(
                self._truncate(r.stdout) or "(no matches)",
                True, time.time() - t0, r.returncode,
            )
        except subprocess.TimeoutExpired:
            self.result = CommandResult("(command timed out)", False, time.time() - t0)
        except Exception as e:
            self.result = CommandResult(f"Error: {e}", False, time.time() - t0)
        return self.result


class ListFilesCommand(ToolCommand):
    name = "list_files"

    def execute(self) -> CommandResult:
        t0 = time.time()
        try:
            path = self.args.get("path", ".")
            if not os.path.isabs(path):
                path = os.path.join(os.getcwd(), path)
            if self.args.get("recursive"):
                r = subprocess.run(
                    ["find", path, "-maxdepth", "3", "-not", "-path", "*/.*"],
                    capture_output=True, text=True, timeout=15,
                )
            else:
                r = subprocess.run(
                    ["ls", "-lah", path],
                    capture_output=True, text=True, timeout=15,
                )
            self.result = CommandResult(
                self._truncate(r.stdout) or "(empty)",
                True, time.time() - t0, r.returncode,
            )
        except subprocess.TimeoutExpired:
            self.result = CommandResult("(command timed out)", False, time.time() - t0)
        except Exception as e:
            self.result = CommandResult(f"Error: {e}", False, time.time() - t0)
        return self.result


class SearchExamplesCommand(ToolCommand):
    """Query DuckDB for successful tool chains similar to a given task description."""
    name = "search_examples"

    def __init__(self, args: dict, storage: "Storage"):
        super().__init__(args)
        self._storage = storage

    def execute(self) -> CommandResult:
        t0 = time.time()
        query = self.args.get("query", "").strip()
        limit = int(self.args.get("limit", 3))
        if not query:
            self.result = CommandResult("Error: query is required", False, time.time() - t0)
            return self.result

        try:
            chains = self._storage.find_similar_chains(query, limit=limit)
        except Exception as e:
            self.result = CommandResult(f"Error: {e}", False, time.time() - t0)
            return self.result

        if not chains:
            self.result = CommandResult(
                "No similar completed tasks found in history.", True, time.time() - t0
            )
            return self.result

        lines = []
        for i, chain in enumerate(chains, 1):
            lines.append(f"--- Example {i} (score={chain['score']}) ---")
            lines.append(f"Task: {chain['description']}")
            lines.append("Steps:")
            for step in chain["steps"]:
                # Redact long args to keep context lean
                args_repr = json.dumps(step["args"])
                if len(args_repr) > 120:
                    args_repr = args_repr[:117] + "..."
                lines.append(f"  {step['tool']}({args_repr})")
                if step["output_summary"]:
                    lines.append(f"    → {step['output_summary'][:100]}")
            lines.append("")

        self.result = CommandResult("\n".join(lines), True, time.time() - t0)
        return self.result


# ── Registry ──────────────────────────────────────────────────────────────────

_COMMAND_MAP: dict[str, type] = {
    "run_bash":        RunBashCommand,
    "run_python":      RunPythonCommand,
    "read_file":       ReadFileCommand,
    "write_file":      WriteFileCommand,
    "grep_search":     GrepSearchCommand,
    "list_files":      ListFilesCommand,
    "search_examples": SearchExamplesCommand,
}


class CommandRegistry:
    def __init__(self, work_dir: str, storage: "Storage"):
        self.work_dir = work_dir
        self.storage = storage

    def build(self, tool_call: dict) -> ToolCommand:
        name = tool_call["function"]["name"]
        args = tool_call["function"].get("arguments", {})
        cls = _COMMAND_MAP.get(name)
        if cls is None:
            raise ValueError(f"Unknown tool: {name}")
        if cls is RunPythonCommand:
            return cls(args, self.work_dir)
        if cls is SearchExamplesCommand:
            return cls(args, self.storage)
        return cls(args)

    def execute_with_cache(
        self, cmd: ToolCommand, task_id: str, session_id: str
    ) -> "CommandResult":
        """Check cache, execute if miss, record result."""
        cmd.timestamp = time.time()

        # Cache check for read-only tools
        cached_output = self.storage.lookup_cached_tool_call(cmd.name, cmd.args)
        if cached_output is not None:
            cmd.result = CommandResult(cached_output, True, 0.0)
            # Still record the cache hit so the graph stays consistent
            self.storage.record_tool_call(task_id, session_id, cmd)
            return cmd.result

        # Execute
        cmd.execute()

        # Persist
        tc_id = self.storage.record_tool_call(task_id, session_id, cmd)

        # Extract entities for the graph (best-effort)
        if cmd.result and tc_id >= 0:
            self.storage.record_entities_from_tool_call(tc_id, cmd.args, cmd.result.output)

        return cmd.result or CommandResult("(no result)", False, 0.0)
